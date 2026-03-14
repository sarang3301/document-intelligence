import os
from pathlib import Path

def load_document(file_path):
    ext = Path(file_path).suffix.lower()
    
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    else:
        return ""

def load_pdf(file_path):
    import fitz  # pymupdf
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    return text

def load_docx(file_path):
    from docx import Document
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def load_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text += para.text + "\n"
    return text